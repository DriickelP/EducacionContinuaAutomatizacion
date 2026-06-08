import sys
import os
import json
import time
import threading
import queue
import logging
import subprocess
import smtplib
import mimetypes
import unicodedata
import re
from datetime import datetime
from pathlib import Path
from email.message import EmailMessage

# --- NUEVAS LIBRERÍAS PARA VARIABLES DE ENTORNO ---
from dotenv import load_dotenv

# Cargar variables del archivo .env
load_dotenv()

# Librerías de Interfaz
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox

# Librerías de Generación y Nube
import openpyxl
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from num2words import num2words as nw
from google.oauth2 import service_account
from googleapiclient.discovery import build

import locale
try:
    locale.setlocale(locale.LC_TIME, "Spanish_Mexico.1252")
except:
    pass
FECHA_ACTUAL = datetime.now().strftime("%B %Y")

try:
    if sys.platform == "win32":
        import win32com.client
        import pythoncom
    EXCEL_ADAPTER_DISPONIBLE = True
except ImportError:
    EXCEL_ADAPTER_DISPONIBLE = False

# ─────────────────────────────────────────────────────────────────────
# CREDENCIALES DE CORREO (Asegúrate de llenarlas con tus datos)
# ─────────────────────────────────────────────────────────────────────
MI_CORREO = os.getenv("CORREO_EMISOR")  
CONTRASENA_APP = os.getenv("CONTRASENA_APP")

# Verificación de seguridad básica (opcional, para avisar en consola si falta el .env)
if not MI_CORREO or not CONTRASENA_APP:
    print("ADVERTENCIA: No se encontraron las credenciales de correo en el archivo .env")

# ─────────────────────────────────────────────────────────────────────
# RUTA BASE Y CARPETAS
# ─────────────────────────────────────────────────────────────────────
def get_base_dir() -> Path:
    if getattr(sys, 'frozen', False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent

BASE_DIR = get_base_dir()
CONFIG_FILE = BASE_DIR / "config.json"
RUTA_BASES_DATOS = BASE_DIR / "bases de datos"
RUTA_MODULOS = BASE_DIR / "modulos"

RUTA_EXCEL = RUTA_BASES_DATOS / "pruebaEC.xlsx"
RUTA_BD = RUTA_BASES_DATOS / "Libro1.xlsx"
CREDENTIALS_FILE = BASE_DIR / "educacion-continua-490016-7ff6875c1d8f.json"

RUTA_IMAGENES = BASE_DIR / "imagenes"
RUTA_RECORTES = RUTA_IMAGENES / "recortes"
RUTA_PLANTILLAS_NUEVAS = BASE_DIR / "PlantillasAct"
RUTA_PLANTILLAS_VIEJAS = BASE_DIR / "PlantillasActIngJorge"
RUTA_SALIDA = BASE_DIR / "Tramites Digitales"

# EL ID DE SEGUIMIENTO AHORA ES FIJO OTRA VEZ
ID_SEGUIMIENTO_NUBE = "1szQLci5kxO10bTGyzQ1wQrQ9z_hITA0uAUHEZ-loaJU"
HOJA_TEC = "SEGUIMIENTO TEC"
HOJA_ESP = "SEGUIMIENTO ESP"
HOJA_BD = "tramites Hechos"

sys.path.append(str(RUTA_MODULOS))
try:
    from face_processor import process_image # type: ignore
    from foto_downloader import descargar_foto # type: ignore
    FACE_PROCESSOR_DISPONIBLE = True
except ImportError:
    FACE_PROCESSOR_DISPONIBLE = False

# ─────────────────────────────────────────────────────────────────────
# ESTILOS GUI
# ─────────────────────────────────────────────────────────────────────
DARK_BG, PANEL_BG, BORDER = "#0f1117", "#1a1d27", "#2a2d3a"
ACCENT, ERROR_COLOR, SUCCESS, WARNING = "#00d4aa", "#ff4757", "#2ed573", "#f5a623"
TEXT_PRIMARY, TEXT_DIM = "#e8eaf0", "#6b7280"
FONT_MONO, FONT_SMALL = ("Consolas", 10), ("Segoe UI", 9)

# ─────────────────────────────────────────────────────────────────────
# UTILIDADES
# ─────────────────────────────────────────────────────────────────────
log_queue = queue.Queue()

class QueueHandler(logging.Handler):
    def emit(self, record):
        log_queue.put(("log", self.format(record), record.levelname))

def matar_proceso(nombre_exe: str):
    if sys.platform == "win32":
        subprocess.call(["taskkill", "/f", "/im", nombre_exe], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

def matar_proceso_excel():
    matar_proceso("excel.exe")

def limpiar_valor(valor):
    """Elimina decimales .0 para no mostrar números como flotantes."""
    if valor is None: return ""
    
    if isinstance(valor, float) and valor.is_integer():
        return str(int(valor))
        
    v_str = str(valor).strip()
    if v_str.endswith(".0"): 
        v_str = v_str[:-2]
    return v_str

def normalizar_condicion(valor):
    """Limpia la celda quitando espacios, pasando a MAYÚSCULAS y quitando acentos."""
    if not valor: 
        return ""
    v = str(valor).strip().upper()
    return v.replace("SÍ", "SI").replace("NÓ", "NO")

def cargar_procesadas(hoja_bd) -> set:
    return {hoja_bd.cell(row=r, column=2).value for r in range(2, hoja_bd.max_row + 1) if hoja_bd.cell(row=r, column=2).value}

def obtener_filas_nube(servicio, nombre_hoja):
    """Obtiene las filas de Google Sheets usando el ID FIJO maestro."""
    try:
        res = servicio.spreadsheets().values().get(spreadsheetId=ID_SEGUIMIENTO_NUBE, range=f"'{nombre_hoja}'!B:B").execute()
        return {limpiar_valor(fila[0]): i + 1 for i, fila in enumerate(res.get("values", [])) if fila}
    except Exception as e: 
        logging.getLogger("EC").error(f"Error obteniendo filas nube: {e}")
        return {}

# ─────────────────────────────────────────────────────────────────────
# LÓGICA DE CORREOS
# ─────────────────────────────────────────────────────────────────────
def enviarESP(ruta, correo_destinatario, log):
    if not correo_destinatario or "@" not in correo_destinatario:
        return False
        
    msg = EmailMessage()
    msg['Subject'] = "Diploma - Trámite Digital"
    msg['From'] = MI_CORREO
    msg['To'] = correo_destinatario
    msg.set_content("Buen día.\n\nEsperando se encuentre bien, por este medio le adjuntamos el diploma correspondiente por haber concluido satisfactoriamente su capacitación.\nAgradecemos la confianza brindada a Educación Continua y le deseamos mucho éxito.\n\nSaludos cordiales.")

    try:
        nombre_archivo = os.path.basename(ruta)
        tipo_mime, _ = mimetypes.guess_type(ruta)
        tipo_principal, sub_tipo = (tipo_mime or 'application/octet-stream').split('/', 1)

        with open(ruta, 'rb') as archivo:
            msg.add_attachment(archivo.read(), maintype=tipo_principal, subtype=sub_tipo, filename=nombre_archivo)
    except FileNotFoundError:
        log.error(f"❌ Error ESP: No se encontró el archivo {ruta} para enviar.")
        return False

    try:
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as servidor:
            servidor.login(MI_CORREO, CONTRASENA_APP)
            servidor.send_message(msg)
        log.info(f"📧 Correo ESP enviado exitosamente a {correo_destinatario}.")
        return True
    except Exception as e:
        log.error(f"❌ Error al enviar correo ESP: {e}")
        return False

def enviarTEC(ruta1, ruta2, correo_destinatario, log):
    if not correo_destinatario or "@" not in correo_destinatario:
        return False

    msg = EmailMessage()
    msg['Subject'] = "Kardex y Diploma - Trámite Digital"
    msg['From'] = MI_CORREO
    msg['To'] = correo_destinatario
    msg.set_content("Buen día.\n\nEsperando se encuentre bien, por este medio le adjuntamos el diploma correspondiente por haber concluido satisfactoriamente su capacitación.\nAgradecemos la confianza brindada a Educación Continua y le deseamos mucho éxito.\n\nSaludos cordiales.")

    for ruta in [ruta1, ruta2]:
        try:
            nombre_archivo = os.path.basename(ruta)
            tipo_mime, _ = mimetypes.guess_type(ruta)
            tipo_principal, sub_tipo = (tipo_mime or 'application/octet-stream').split('/', 1)

            with open(ruta, 'rb') as archivo:
                msg.add_attachment(archivo.read(), maintype=tipo_principal, subtype=sub_tipo, filename=nombre_archivo)
        except FileNotFoundError:
            log.error(f"❌ Error TEC: No se encontró el archivo {ruta} para enviar.")
            return False

    try:
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as servidor:
            servidor.login(MI_CORREO, CONTRASENA_APP)
            servidor.send_message(msg)
        log.info(f"📧 Correo TEC (2 archivos) enviado exitosamente a {correo_destinatario}.")
        return True
    except Exception as e:
        log.error(f"❌ Error al enviar correo TEC: {e}")
        return False

# ─────────────────────────────────────────────────────────────────────
# LÓGICA DE GENERACIÓN (Kardex y Diplomas)
# ─────────────────────────────────────────────────────────────────────
def limpiar_para_ruta(texto: str) -> str:
    """Elimina acentos, caracteres inválidos de Windows y espacios al final para evitar errores de COM."""
    if not texto: return "Desconocido"
    texto_limpio = "".join(c for c in unicodedata.normalize('NFD', str(texto)) if unicodedata.category(c) != 'Mn')
    texto_limpio = re.sub(r'[<>:"/\\|?*]', '', texto_limpio)
    return texto_limpio.strip()

def procesar_imagen(nombre: str, matricula: str, log) -> tuple:
    if not FACE_PROCESSOR_DISPONIBLE: return None, False
    for ext in (".jpg", ".jpeg", ".png", ".webp"):
        ruta_img = RUTA_IMAGENES / f"{nombre}{ext}"
        if ruta_img.exists(): break
    else:
        return None, False

    ruta_recorte = RUTA_RECORTES / f"{nombre}.png"
    try:
        process_image(input_path=str(ruta_img), output_path=str(ruta_recorte), mode="local")
        return ruta_recorte, True
    except Exception as e:
        log.error(f"Error en foto de {nombre}: {e}")
        return None, False

def generar_diploma(nombre: str, carrera: str, ppt_app, log, ruta_plantillas):
    carrera_ruta = limpiar_para_ruta(carrera)
    nombre_ruta = limpiar_para_ruta(nombre)
    
    carpeta_salida = RUTA_SALIDA / carrera_ruta
    carpeta_salida.mkdir(parents=True, exist_ok=True)
    plantilla = ruta_plantillas / f"D{carrera.strip()}.pptx"
    
    if not plantilla.exists():
        log.error(f"Plantilla no encontrada: {plantilla}")
        return None

    prs = Presentation(str(plantilla))
    
    colecciones = [
        prs.slides[0].shapes,
        prs.slides[0].slide_layout.shapes,
        prs.slides[0].slide_layout.slide_master.shapes
    ]

    for shapes in colecciones:
        for shape in shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                texto_completo = "".join(r.text for r in para.runs)
                
                if "##nombre##" in texto_completo or "##fecha##" in texto_completo or "{{fecha}}" in texto_completo:
                    es_nombre = "##nombre##" in texto_completo
                    
                    nuevo = texto_completo.replace("##nombre##", nombre.title())
                    nuevo = nuevo.replace("##fecha##", FECHA_ACTUAL)
                    nuevo = nuevo.replace("{{fecha}}", FECHA_ACTUAL)
                    
                    for idx, run in enumerate(para.runs):
                        run.text = nuevo if idx == 0 else ""
                        if idx == 0 and es_nombre:
                            run.font.name = "Arial"
                            run.font.size = Pt(32)
                            run.font.bold = True
                            
                    if es_nombre:
                        para.alignment = PP_ALIGN.CENTER
                        
            shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    ruta_pptx = str((carpeta_salida / f"{nombre_ruta}.pptx").resolve())
    ruta_pdf = ruta_pptx.replace(".pptx", ".pdf")
    prs.save(ruta_pptx)

    if ppt_app:
        pres = ppt_app.Presentations.Open(ruta_pptx, WithWindow=False)
        pres.SaveAs(ruta_pdf, 32)
        pres.Close()
        try: os.remove(ruta_pptx)
        except: pass
        return ruta_pdf
    return ruta_pptx

def generar_kardex(nombre: str, matricula: str, carrera: str, califs: list, ruta_foto, word_app, log, ruta_plantillas):
    carrera_ruta = limpiar_para_ruta(carrera)
    nombre_ruta = limpiar_para_ruta(nombre)
    
    carpeta_salida = RUTA_SALIDA / carrera_ruta
    carpeta_salida.mkdir(parents=True, exist_ok=True)
    plantilla = ruta_plantillas / f"{carrera.strip()}.docx"
    
    if not plantilla.exists():
        log.error(f"Plantilla kardex no encontrada: {plantilla}")
        return None

    reemplazos = {"{{nombre}}": nombre.title(), "{{matricula}}": str(matricula), "{{fecha}}": FECHA_ACTUAL}
    
    for i, cal in enumerate(califs):
        if isinstance(cal, float) and cal.is_integer():
            cal = int(cal)
            
        c_val = cal if cal is not None else 0
        reemplazos[f"{{{{cal{i+1}}}}}"] = str(c_val)
        
        try: 
            num_limpio = int(float(c_val)) if str(c_val).replace('.', '', 1).isdigit() else c_val
            reemplazos[f"{{{{num{i+1}}}}}"] = nw(num_limpio, lang="es").upper()
        except: 
            reemplazos[f"{{{{num{i+1}}}}}"] = str(c_val)

    doc = Document(str(plantilla))
    for p in doc.paragraphs:
        for run in p.runs:
            for c, v in reemplazos.items():
                if c in run.text: run.text = run.text.replace(c, v)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    texto = p.text
                    modificado = False
                    for c, v in reemplazos.items():
                        if c in texto:
                            texto = texto.replace(c, v)
                            modificado = True
                    if modificado:
                        p.clear()
                        run = p.add_run(texto)
                        run.font.name = "Switzerland"
                        run.font.size = Pt(10)

    ruta_docx = str((carpeta_salida / f"{nombre_ruta.title()} KARDEX.docx").resolve())
    ruta_pdf = ruta_docx.replace(".docx", ".pdf")
    doc.save(ruta_docx)

    if word_app:
        doc_word = word_app.Documents.Open(ruta_docx)
        if ruta_foto and ruta_foto.exists():
            rango = doc_word.Paragraphs(1).Range
            rango.Collapse(0)
            shape = doc_word.InlineShapes.AddPicture(FileName=str(ruta_foto), LinkToFile=False, SaveWithDocument=True, Range=rango)
            floating = shape.ConvertToShape()
            floating.ZOrder(5)
            floating.LockAspectRatio = True
            floating.Height = 3.57 * 28.35
            floating.Left = 14.5 * 28.35
            floating.Top = 1.3 * 28.35
            floating.Left = floating.Left + (((3.21 * 28.35) - floating.Width) / 2)

        doc_word.SaveAs(ruta_pdf, FileFormat=17)
        doc_word.Close(SaveChanges=False)
        try: os.remove(ruta_docx)
        except: pass
        return ruta_pdf
    return ruta_docx

def generar_diploma_esp(nombre: str, carrera: str, horas: str, ppt_app, log, ruta_plantillas):
    carrera_ruta = limpiar_para_ruta(carrera)
    nombre_ruta = limpiar_para_ruta(nombre)
    
    carpeta_salida = RUTA_SALIDA / carrera_ruta
    carpeta_salida.mkdir(parents=True, exist_ok=True)
    plantilla = ruta_plantillas / "DIPLOMADOS.pptx"
    
    if not plantilla.exists():
        log.error(f"Plantilla ESP no encontrada: {plantilla}")
        return None


    prs = Presentation(str(plantilla))
    

    try:
        h_float = float(horas)
        h_val = int(h_float)
        horas_limpias = str(h_val) if h_float.is_integer() else str(h_float)
    except ValueError:
        h_val = 0
        horas_limpias = str(horas)

    t_curso = "Diplomado" if h_val >= 40 else "Curso"
    conector = "en" if h_val >= 40 else "de"

    reemplazos = {
        "##nombre##": nombre.title(), "##fecha##": FECHA_ACTUAL,
        "##diplomado##": carrera.strip(), "##hr##": horas_limpias,
        "##tipocurso##": t_curso, "##conector##": conector
    }

    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    def reemplazar_xml(elemento):
        for p_elem in elemento.iter(f"{{{NS}}}p"):
            runs = p_elem.findall(f"{{{NS}}}r")
            if not runs: continue
            texto = "".join((r.find(f"{{{NS}}}t").text or "") for r in runs if r.find(f"{{{NS}}}t") is not None)
            if not any(k in texto for k in reemplazos): continue

            nuevo = texto
            for c, v in reemplazos.items(): nuevo = nuevo.replace(c, v)

            primer_t = runs[0].find(f"{{{NS}}}t")
            if primer_t is not None: primer_t.text = nuevo
            for run in runs[1:]:
                t = run.find(f"{{{NS}}}t")
                if t is not None: t.text = ""

            rPr = runs[0].find(f"{{{NS}}}rPr")
            if rPr is None:
                from lxml import etree
                rPr = etree.SubElement(runs[0], f"{{{NS}}}rPr")
                runs[0].insert(0, rPr)

            if "##nombre##" in texto:
                rPr.set("sz", "2400"); rPr.set("b", "1")
            elif "##diplomado##" in texto:
                rPr.set("sz", "1800"); rPr.set("b", "1")

    reemplazar_xml(prs.slides[0]._element)
    reemplazar_xml(prs.slides[0].slide_layout.slide_master._element)

    ruta_pptx = str((carpeta_salida / f"{nombre_ruta}.pptx").resolve())
    ruta_pdf = ruta_pptx.replace(".pptx", ".pdf")
    prs.save(ruta_pptx)

    if ppt_app:
        pres = ppt_app.Presentations.Open(ruta_pptx, WithWindow=False)
        pres.SaveAs(ruta_pdf, 32)
        pres.Close()
        try: os.remove(ruta_pptx)
        except: pass
        return ruta_pdf
    return ruta_pptx

# ─────────────────────────────────────────────────────────────────────
# GUI & ORQUESTADOR PRINCIPAL
# ─────────────────────────────────────────────────────────────────────
class SettingsWindow(tk.Toplevel):
    def __init__(self, parent):
        super().__init__(parent)
        self.title("Ajustes")
        self.geometry("500x420")
        self.configure(bg=DARK_BG)
        self.transient(parent)
        self.grab_set()
        self.parent = parent
        self._build_ui()

    def _build_ui(self):
        tk.Label(self, text="CONFIGURAR FORMULARIOS", font=("Segoe UI", 12, "bold"), bg=DARK_BG, fg=ACCENT).pack(pady=15)
        f = tk.Frame(self, bg=DARK_BG); f.pack(fill="x", padx=20)
        self.entry_id = tk.Entry(f, bg=PANEL_BG, fg=TEXT_PRIMARY, insertbackground=ACCENT, font=FONT_MONO, bd=0)
        self.entry_id.pack(side="left", fill="x", expand=True, padx=(0, 10), ipady=5)
        tk.Button(f, text="Añadir", bg=ACCENT, fg=DARK_BG, font=FONT_SMALL, relief="flat", padx=15, command=self._add).pack(side="right")

        lf = tk.Frame(self, bg=PANEL_BG); lf.pack(fill="both", expand=True, padx=20, pady=(15, 15))
        self.lb = tk.Listbox(lf, bg=PANEL_BG, fg=TEXT_PRIMARY, font=FONT_MONO, bd=0, selectbackground=ACCENT)
        self.lb.pack(side="left", fill="both", expand=True, padx=5, pady=5)
        for i in self.parent.config.get("forms_ids", []): self.lb.insert(tk.END, i)
        tk.Button(self, text="Eliminar", bg=ERROR_COLOR, fg="white", font=FONT_SMALL, relief="flat", command=self._del).pack(pady=10)

    def _add(self):
        raw = self.entry_id.get().strip()
        if raw:
            cid = raw.split("spreadsheets/d/")[1].split("/")[0] if "spreadsheets/d/" in raw else raw
            if cid not in self.parent.config["forms_ids"]:
                self.parent.config["forms_ids"].append(cid); self.lb.insert(tk.END, cid)
                self.entry_id.delete(0, tk.END); self.parent._save_config()

    def _del(self):
        s = self.lb.curselection()
        if s:
            self.parent.config["forms_ids"].remove(self.lb.get(s[0]))
            self.lb.delete(s[0]); self.parent._save_config()


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Trámites Digitales FIME - Final")
        self.geometry("1100x750") 
        self.configure(bg=DARK_BG)
        self.config = self._load_config()

        # Variables de Control
        self.running = False
        self.is_paused = False
        self.pause_event = threading.Event()
        self.pause_event.set()

        # Variables del Dashboard
        self.tot_tec = 0
        self.tot_esp = 0
        self.tot_correos = 0
        self.tot_omitidos = 0
        self.errores = 0
        
        self._build_ui(); self._setup_log(); self._poll()

    def _load_config(self): 
        if CONFIG_FILE.exists():
            return json.load(open(CONFIG_FILE))
        return {"forms_ids": []}
        
    def _save_config(self): json.dump(self.config, open(CONFIG_FILE, "w"), indent=4)
    def _setup_log(self):
        log = logging.getLogger("EC"); log.setLevel(logging.INFO)
        if not log.handlers:
            h = QueueHandler()
            h.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s", "%H:%M:%S"))
            log.addHandler(h)

    def _build_ui(self):
        h = tk.Frame(self, bg=DARK_BG); h.pack(fill="x", padx=24, pady=(20, 0))
        tk.Label(h, text="TRAMITES DIGITALES", font=("Segoe UI", 22, "bold"), bg=DARK_BG, fg=ACCENT).pack(side="left")

        self.modo_prueba = tk.BooleanVar(value=False)
        frm_toggle = tk.Frame(h, bg=DARK_BG)
        frm_toggle.pack(side="right", padx=(10, 0))
        tk.Label(frm_toggle, text="MODO PRUEBA", font=("Segoe UI", 8, "bold"), bg=DARK_BG, fg=WARNING).pack(side="left", padx=(0, 6))
        self.chk_prueba = tk.Checkbutton(
            frm_toggle, variable=self.modo_prueba,
            bg=DARK_BG, activebackground=DARK_BG,
            fg=WARNING, selectcolor=PANEL_BG,
            command=self._toggle_modo_prueba
        )
        self.chk_prueba.pack(side="left")

        tk.Button(h, text="⚙️ Ajustes", font=FONT_SMALL, bg=PANEL_BG, fg=TEXT_PRIMARY, relief="flat", padx=10, command=lambda: SettingsWindow(self)).pack(side="right")
        tk.Frame(self, bg=BORDER, height=1).pack(fill="x", padx=24, pady=12)

        body = tk.Frame(self, bg=DARK_BG); body.pack(fill="both", expand=True, padx=24)
        body.columnconfigure(0, weight=3); body.columnconfigure(1, weight=1); body.rowconfigure(0, weight=1)

        style = ttk.Style(); style.theme_use("clam")
        style.configure("EC.TNotebook", background=PANEL_BG, borderwidth=0)
        style.configure("EC.TNotebook.Tab", background=DARK_BG, foreground=TEXT_DIM, padding=[15, 6], font=("Segoe UI", 9, "bold"))
        style.map("EC.TNotebook.Tab", background=[("selected", PANEL_BG)], foreground=[("selected", ACCENT)])

        self.nb = ttk.Notebook(body, style="EC.TNotebook"); self.nb.grid(row=0, column=0, sticky="nsew", padx=(0, 15))

        self.tab_res = tk.Frame(self.nb, bg=DARK_BG)
        self.nb.add(self.tab_res, text=" Resumen Ejecutivo ")
        self.res_box = scrolledtext.ScrolledText(self.tab_res, bg=DARK_BG, fg=TEXT_PRIMARY, font=FONT_MONO, bd=0, state="disabled")
        self.res_box.pack(fill="both", expand=True, padx=10, pady=10)

        self.tab_con = tk.Frame(self.nb, bg=DARK_BG)
        self.nb.add(self.tab_con, text=" Log de Sistema ")
        self.con_box = scrolledtext.ScrolledText(self.tab_con, bg=DARK_BG, fg=TEXT_PRIMARY, font=FONT_MONO, bd=0, state="disabled")
        self.con_box.pack(fill="both", expand=True, padx=10, pady=10)

        for b in [self.res_box, self.con_box]:
            b.tag_config("TEC", foreground=ACCENT, font=("Consolas", 10, "bold"))
            b.tag_config("ESP", foreground=SUCCESS, font=("Consolas", 10, "bold"))
            b.tag_config("INFO", foreground=TEXT_PRIMARY)
            b.tag_config("ERROR", foreground=ERROR_COLOR)
            b.tag_config("WARNING", foreground=WARNING)
            b.tag_config("SUCCESS", foreground=SUCCESS)

        # NUEVO DASHBOARD (LADO DERECHO)
        right = tk.Frame(body, bg=DARK_BG); right.grid(row=0, column=1, sticky="nsew")
        
        # 1. Cuadrícula de Contadores
        frame_counters = tk.Frame(right, bg=DARK_BG)
        frame_counters.pack(fill="x", pady=(0, 10))
        frame_counters.columnconfigure(0, weight=1)
        frame_counters.columnconfigure(1, weight=1)

        self.lbl_tec = self._card_grid(frame_counters, "🎓 TEC", "0", ACCENT, 0, 0)
        self.lbl_esp = self._card_grid(frame_counters, "📜 ESP", "0", "#a29bfe", 0, 1)
        self.lbl_correos = self._card_grid(frame_counters, "📧 CORREOS", "0", "#74b9ff", 1, 0)
        self.lbl_omit = self._card_grid(frame_counters, "⏭️ OMITIDOS", "0", WARNING, 1, 1)
        self.lbl_err = self._card_grid(frame_counters, "❌ ERRORES", "0", ERROR_COLOR, 2, 0, columnspan=2)

        # 2. Tarjeta Dinámica de "Último Procesado"
        self.frame_last = tk.Frame(right, bg=PANEL_BG, highlightbackground=BORDER, highlightthickness=1)
        self.frame_last.pack(fill="both", expand=True, pady=(10, 10))
        
        tk.Label(self.frame_last, text="ÚLTIMO TRÁMITE PROCESADO", font=("Segoe UI", 9, "bold"), bg=PANEL_BG, fg=TEXT_DIM).pack(anchor="w", padx=12, pady=(15, 10))
        
        self.lbl_last_name = tk.Label(self.frame_last, text="Esperando ejecución...", font=("Segoe UI", 12, "bold"), bg=PANEL_BG, fg=TEXT_PRIMARY, wraplength=250, justify="left")
        self.lbl_last_name.pack(anchor="w", padx=12)
        
        self.lbl_last_details = tk.Label(self.frame_last, text="\n\n", font=("Segoe UI", 10), bg=PANEL_BG, fg=TEXT_DIM, justify="left")
        self.lbl_last_details.pack(anchor="w", padx=12, pady=(8, 15))

        # Botón de Inicio
        self.btn = tk.Button(right, text="▶  INICIAR", font=("Segoe UI", 12, "bold"), bg=ACCENT, fg=DARK_BG, relief="flat", pady=12, command=self._start)
        self.btn.pack(fill="x", side="bottom")

        footer = tk.Frame(self, bg=DARK_BG)
        footer.pack(fill="x", padx=24, pady=10)
        self.status_label = tk.Label(footer, text="● En espera", font=FONT_SMALL, bg=DARK_BG, fg=TEXT_DIM)
        self.status_label.pack(side="left")

    def _card_grid(self, parent, title, val, color, r, c, columnspan=1):
        frm = tk.Frame(parent, bg=PANEL_BG, highlightbackground=BORDER, highlightthickness=1)
        frm.grid(row=r, column=c, columnspan=columnspan, sticky="nsew", padx=4, pady=4)
        tk.Label(frm, text=title, font=("Segoe UI", 8, "bold"), bg=PANEL_BG, fg=TEXT_DIM).pack(anchor="w", padx=10, pady=(8, 0))
        l = tk.Label(frm, text=val, font=("Segoe UI", 22, "bold"), bg=PANEL_BG, fg=color)
        l.pack(anchor="w", padx=10, pady=(0, 8))
        return l

    def _toggle_modo_prueba(self):
        if self.modo_prueba.get():
            self.status_label.config(text="● MODO PRUEBA ACTIVO — Sin escritura en Sheets ni envío de correo", fg=WARNING)
            log_queue.put(("log", "⚠️ MODO PRUEBA ACTIVADO: No se escribirá en Sheets ni se enviarán correos reales.", "WARNING"))
        else:
            self.status_label.config(text="● En espera", fg=TEXT_DIM)
            log_queue.put(("log", "✅ MODO PRUEBA DESACTIVADO: Operación normal.", "INFO"))

    def _start(self):
        if not self.config.get("forms_ids"): 
            return messagebox.showerror("Error", "Añade al menos un ID de formulario en Ajustes.")
            
        if self.running: return

        self.running = True
        self.is_paused = False
        self.pause_event.set()

        # Reseteo del Dashboard
        self.tot_tec = 0; self.tot_esp = 0; self.tot_correos = 0; self.tot_omitidos = 0; self.errores = 0
        self.lbl_tec.config(text="0"); self.lbl_esp.config(text="0"); self.lbl_correos.config(text="0")
        self.lbl_omit.config(text="0"); self.lbl_err.config(text="0")
        self.lbl_last_name.config(text="Iniciando...", fg=ACCENT)
        self.lbl_last_details.config(text="Conectando a los servicios...\nPor favor espera.")
        
        self.matriculas_vistas = set()

        self.btn.config(text="⏸  PAUSAR", bg=WARNING, command=self._toggle_pause)
        self.status_label.config(text="● Procesando...", fg=ACCENT)

        self.res_box.config(state="normal"); self.res_box.delete("1.0", "end"); self.res_box.config(state="disabled")
        self.con_box.config(state="normal"); self.con_box.delete("1.0", "end"); self.con_box.config(state="disabled")
        self.nb.select(self.tab_res)

        threading.Thread(target=self._worker, daemon=True).start()

    def _toggle_pause(self):
        if self.is_paused:
            self.is_paused = False
            self.pause_event.set()
            self.btn.config(text="⏸  PAUSAR", bg=WARNING)
            self.status_label.config(text="● Procesando...", fg=ACCENT)
            log_queue.put(("log", "▶ Proceso reanudado.", "INFO"))
        else:
            self.is_paused = True
            self.pause_event.clear()
            self.btn.config(text="▶  REANUDAR", bg=SUCCESS)
            self.status_label.config(text="● Pausado", fg=WARNING)
            log_queue.put(("log", "⏸ Proceso pausado por el usuario.", "WARNING"))

    def _worker(self):
        log = logging.getLogger("EC")
        word_app, ppt_app, excel = None, None, None
        modo_prueba = self.modo_prueba.get()

        if sys.platform == "win32": pythoncom.CoInitialize()

        if not RUTA_EXCEL.exists() or not RUTA_BD.exists():
            log_queue.put(("error", f"❌ Falta archivo Excel o Libro1.xlsx en 'bases de datos'", "ERROR"))
            if sys.platform == "win32": pythoncom.CoUninitialize()
            return

        try:
            log.info("Conectando con Google Forms y Sheets...")
            creds = service_account.Credentials.from_service_account_file(str(CREDENTIALS_FILE), scopes=["https://www.googleapis.com/auth/spreadsheets"])
            serv = build("sheets", "v4", credentials=creds, cache_discovery=False)

            m_tec, m_esp = {}, {}
            matriculas_en_formularios = set() 
            
            for fid in self.config.get("forms_ids", []):
                try:
                    res = serv.spreadsheets().values().get(spreadsheetId=fid, range="'Respuestas de formulario 1'!A:AK").execute()
                    for fila in res.get("values", []):
                        
                        correo_form = limpiar_valor(fila[35]) if len(fila) > 35 else ""
                        
                        mat_tec = limpiar_valor(fila[7]) if len(fila) > 7 else ""
                        if mat_tec: 
                            m_tec[mat_tec] = {
                                "nombre": str(fila[8]).strip().title() if len(fila) > 8 else "Desconocido",
                                "correo": correo_form
                            }
                            matriculas_en_formularios.add(mat_tec)
                            
                        mat_esp = limpiar_valor(fila[3]) if len(fila) > 3 else ""
                        if mat_esp: 
                            m_esp[mat_esp] = {
                                "nombre": str(fila[4]).strip().title() if len(fila) > 4 else "Desconocido",
                                "correo": correo_form
                            }
                            matriculas_en_formularios.add(mat_esp)
                except Exception as e: 
                    log.warning(f"No se pudo leer el formulario con ID: {fid}. Error: {e}")

            filas_tec = obtener_filas_nube(serv, HOJA_TEC)
            filas_esp = obtener_filas_nube(serv, HOJA_ESP)
            lote_nube = []

            self.pause_event.wait()

            log.info("Inicializando motores de Office...")
            excel = win32com.client.DispatchEx("Excel.Application")
            excel.Visible = False; excel.DisplayAlerts = False
            word_app = win32com.client.DispatchEx("Word.Application")
            word_app.Visible = False; word_app.DisplayAlerts = False
            ppt_app = win32com.client.DispatchEx("PowerPoint.Application")

            wb_bd = openpyxl.load_workbook(str(RUTA_BD))
            h_bd = wb_bd[HOJA_BD]
            procesadas = cargar_procesadas(h_bd)
            sig_fila_bd = h_bd.max_row + 1

            log.info("🔄 Refrescando conexiones web de Seguimiento...")
            wb = excel.Workbooks.Open(os.path.abspath(RUTA_EXCEL))

            for c in wb.Connections:
                try: c.OLEDBConnection.BackgroundQuery = False
                except: pass
                try: c.ODBCConnection.BackgroundQuery = False
                except: pass
                try: c.DataFeedConnection.BackgroundQuery = False  
                except: pass

            wb.RefreshAll()
            
            log.info("⏳ Esperando a que Excel termine de descargar los datos de la nube...")
            tiempo_espera = 0
            while tiempo_espera < 60: 
                pythoncom.PumpWaitingMessages() 
                
                refrescando = False
                for c in wb.Connections:
                    try: 
                        if c.OLEDBConnection.Refreshing: refrescando = True
                    except: pass
                    try: 
                        if c.ODBCConnection.Refreshing: refrescando = True
                    except: pass
                
                if not refrescando and excel.CalculationState == 0:
                    break
                    
                time.sleep(1)
                tiempo_espera += 1

            excel.CalculateUntilAsyncQueriesDone()
            wb.Save() 
            log.info("✅ Descarga de la nube y refresco completados con éxito.")

            if modo_prueba:
                log.info("🧪 MODO PRUEBA ACTIVO: Actuando como Diodo. Se bloquea la escritura en Sheets, Excel y Correos.")

            # ── FLUJO TEC ─────────────────────────────────────────────
            log.info("━" * 30 + "\nFLUJO TEC (Kardex, Diploma y Correo)")
            try:
                h_tec = wb.Sheets(HOJA_TEC)
                max_r = h_tec.Cells(h_tec.Rows.Count, 2).End(-4162).Row
                if max_r >= 2:
                    dt = h_tec.Range(h_tec.Cells(1, 1), h_tec.Cells(max_r, 36)).Value
                    for f in range(2, max_r + 1):
                        self.pause_event.wait()
                        if not self.running: break

                        fila = dt[f - 1]
                        mat = limpiar_valor(fila[1])
                        carrera = fila[3] or "Desconocida"
                        
                        # --- PLANTILLAS: siempre PlantillasAct ---
                        ruta_plantillas_actual = RUTA_PLANTILLAS_NUEVAS

                        celda_envio = normalizar_condicion(fila[5])
                        celda_papeleria = normalizar_condicion(fila[7])
                        
                        correo_excel = limpiar_valor(fila[35]) if len(fila) > 35 else ""
                        correo_destinatario = ""
                        nombre_actual = str(fila[2]).strip().title() if fila[2] else "Desconocido"

                        if mat in m_tec:
                            nombre_actual = m_tec[mat]["nombre"]
                            correo_destinatario = m_tec[mat]["correo"]
                            if not modo_prueba:
                                h_tec.Cells(f, 3).Value = nombre_actual
                                
                        if not correo_destinatario:
                            correo_destinatario = correo_excel
                
                        if celda_envio == "SI" and celda_papeleria == "NO" and mat not in procesadas:
                            
                            if mat not in matriculas_en_formularios:
                                log.warning(f"⏭️ {nombre_actual} ({mat}) omitido TEC: Cumple en Excel, pero NO está en las respuestas de los formularios.")
                                log_queue.put(("metric", "omitido", "INFO"))
                                continue

                            log.info(f"⚙️ Procesando TEC | Matrícula: {mat} | Oferta: {carrera} | Alumno: {nombre_actual}")
                            
                            ruta_descargada, nombre_descargado = descargar_foto(matricula=mat, carpeta_destino=str(RUTA_IMAGENES))
                            ruta_foto, exito = procesar_imagen(nombre_actual, mat, log)
                            
                            if not exito:
                                if not modo_prueba:
                                    h_tec.Cells(f, 6).Value = "Foto invalida"
                                    if mat in filas_tec:
                                        lote_nube.append({"range": f"'{HOJA_TEC}'!F{filas_tec[mat]}", "values": [["Foto invalida"]]})
                                    log_queue.put(("metric", "omitido", "INFO"))
                                    continue
                                else:
                                    log.warning(f"⚠️ Sin foto para {nombre_actual} — se generará kardex sin imagen.")
                                    ruta_foto = None

                            califs = [fila[c - 1] for c in range(13, 24)]

                            ruta_diploma = generar_diploma(nombre_actual, carrera, ppt_app, log, ruta_plantillas_actual)
                            ruta_kardex = generar_kardex(nombre_actual, mat, carrera, califs, ruta_foto, word_app, log, ruta_plantillas_actual)

                            if not modo_prueba:
                                estado_tramite = "Pendiente"
                                
                                h_bd.cell(row=sig_fila_bd, column=1, value=nombre_actual)
                                h_bd.cell(row=sig_fila_bd, column=2, value=mat)
                                procesadas.add(mat); sig_fila_bd += 1
                                
                                if correo_destinatario and "@" in correo_destinatario:
                                    if ruta_diploma and ruta_kardex:
                                        envio_ok = enviarTEC(ruta_diploma, ruta_kardex, correo_destinatario, log)
                                        if envio_ok: 
                                            log_queue.put(("metric", "correo", "INFO"))
                                            estado_tramite = "Si"
                                            h_tec.Cells(f, 11).Value = "Correo enviado (Por codigo)"
                                            if mat in filas_tec:
                                                lote_nube.append({"range": f"'{HOJA_TEC}'!K{filas_tec[mat]}", "values": [["Correo enviado (Por codigo)"]]})
                                else:
                                    h_tec.Cells(f, 11).Value = "No se encontró su correo"
                                    if mat in filas_tec:
                                        lote_nube.append({"range": f"'{HOJA_TEC}'!K{filas_tec[mat]}", "values": [["No se encontró su correo"]]})
                                    log.warning(f"⚠️ Sin correo válido para {nombre_actual}. Se marca en columna K.")
                                
                                h_tec.Cells(f, 8).Value = estado_tramite
                                if mat in filas_tec:
                                    lote_nube.append({"range": f"'{HOJA_TEC}'!H{filas_tec[mat]}", "values": [[estado_tramite]]})

                            else:
                                if correo_destinatario and "@" in correo_destinatario:
                                    log.info(f"🧪 PRUEBA: Simulación de correo TEC a {correo_destinatario}")
                                else:
                                    log.info(f"🧪 PRUEBA: Simulación de escritura 'No se encontró su correo' en Columna K para {nombre_actual}")

                            log_queue.put(("metric", "tec", "INFO"))
                            
                            msg_resumen = f"[TEC] {'🧪 PRUEBA | ' if modo_prueba else ''}Matrícula: {mat}\n      👤 {nombre_actual}\n      🎓 Oferta: {carrera}\n      ↳ Documentos generados correctamente."
                            log_queue.put(("res_tec", msg_resumen, "SUCCESS", nombre_actual, mat, carrera, ""))

            except Exception as e: log.error(f"Error TEC: {e}")

            # ── FLUJO ESP ─────────────────────────────────────────────
            log.info("━" * 30 + "\nFLUJO ESP (Diplomas y Correo)")
            try:
                h_esp = wb.Sheets(HOJA_ESP)
                max_r = h_esp.Cells(h_esp.Rows.Count, 2).End(-4162).Row
                if max_r >= 2:
                    dt = h_esp.Range(h_esp.Cells(1, 1), h_esp.Cells(max_r, 36)).Value
                    for f in range(2, max_r + 1):
                        self.pause_event.wait()
                        if not self.running: break

                        fila = dt[f - 1]
                        mat = limpiar_valor(fila[1])
                        carrera = fila[4] or "Desconocida"
                        horas = limpiar_valor(fila[3]) or "0"
                        
                        # --- PLANTILLAS: siempre PlantillasAct ---
                        ruta_plantillas_actual = RUTA_PLANTILLAS_NUEVAS
                        
                        celda_envio = normalizar_condicion(fila[7])
                        celda_papeleria = normalizar_condicion(fila[9])
                        
                        correo_excel = limpiar_valor(fila[35]) if len(fila) > 35 else ""
                        correo_destinatario = ""
                        nombre_actual = str(fila[2]).strip().title() if fila[2] else "Desconocido"

                        if mat in m_esp:
                            nombre_actual = m_esp[mat]["nombre"]
                            correo_destinatario = m_esp[mat]["correo"]
                            if not modo_prueba:
                                h_esp.Cells(f, 3).Value = nombre_actual
                                
                        if not correo_destinatario:
                            correo_destinatario = correo_excel

                        if celda_envio == "SI" and celda_papeleria == "NO" and mat not in procesadas:
                            
                            if mat not in matriculas_en_formularios:
                                log.warning(f"⏭️ {nombre_actual} ({mat}) omitido ESP: Cumple en Excel, pero NO está en las respuestas de los formularios.")
                                log_queue.put(("metric", "omitido", "INFO"))
                                continue

                            log.info(f"⚙️ Procesando ESP | Matrícula: {mat} | Oferta: {carrera} ({horas} hrs) | Alumno: {nombre_actual}")
                            
                            ruta_diploma = generar_diploma_esp(nombre_actual, carrera, str(horas), ppt_app, log, ruta_plantillas_actual)

                            if not modo_prueba:
                                estado_tramite = "Pendiente"
                                
                                h_bd.cell(row=sig_fila_bd, column=1, value=nombre_actual)
                                h_bd.cell(row=sig_fila_bd, column=2, value=mat)
                                procesadas.add(mat); sig_fila_bd += 1
                                
                                if correo_destinatario and "@" in correo_destinatario:
                                    if ruta_diploma:
                                        envio_ok = enviarESP(ruta_diploma, correo_destinatario, log)
                                        if envio_ok: 
                                            log_queue.put(("metric", "correo", "INFO"))
                                            estado_tramite = "Si"
                                            h_esp.Cells(f, 13).Value = "Correo enviado (Por codigo)"
                                            if mat in filas_esp:
                                                lote_nube.append({"range": f"'{HOJA_ESP}'!M{filas_esp[mat]}", "values": [["Correo enviado (Por codigo)"]]})
                                else:
                                    h_esp.Cells(f, 13).Value = "No se encontró su correo"
                                    if mat in filas_esp:
                                        lote_nube.append({"range": f"'{HOJA_ESP}'!M{filas_esp[mat]}", "values": [["No se encontró su correo"]]})
                                    log.warning(f"⚠️ Sin correo válido para {nombre_actual}. Se marca en columna M.")
                                
                                h_esp.Cells(f, 10).Value = estado_tramite
                                if mat in filas_esp:
                                    lote_nube.append({"range": f"'{HOJA_ESP}'!J{filas_esp[mat]}", "values": [[estado_tramite]]})

                            else:
                                if correo_destinatario and "@" in correo_destinatario:
                                    log.info(f"🧪 PRUEBA: Simulación de correo ESP a {correo_destinatario}")
                                else:
                                    log.info(f"🧪 PRUEBA: Simulación de escritura 'No se encontró su correo' en Columna M para {nombre_actual}")

                            log_queue.put(("metric", "esp", "INFO"))
                            
                            msg_resumen = f"[ESP] {'🧪 PRUEBA | ' if modo_prueba else ''}Matrícula: {mat}\n      👤 {nombre_actual}\n      📜 Oferta: {carrera} ({horas} hrs)\n      ↳ Diploma generado correctamente."
                            log_queue.put(("res_esp", msg_resumen, "SUCCESS", nombre_actual, mat, carrera, ""))

            except Exception as e: log.error(f"Error ESP: {e}")

            # ── GUARDADO FINAL ────────────────────────────────────────
            if not modo_prueba:
                # batchUpdate — una sola petición para todo el lote
                if lote_nube:
                    try:
                        serv.spreadsheets().values().batchUpdate(
                            spreadsheetId=ID_SEGUIMIENTO_NUBE,
                            body={
                                "valueInputOption": "USER_ENTERED",
                                "data": lote_nube
                            }
                        ).execute()
                        log.info("✅ Nube actualizada: %d celdas en una sola petición", len(lote_nube))
                    except Exception as e:
                        log.error(f"Error actualizando la nube en lote: {e}")

                wb.Save()
                wb_bd.save(str(RUTA_BD))
                log.info("✅ Escritura en Google Sheets y Excel local completada.")
            else:
                log.info("✅ MODO PRUEBA: Fin de ejecución sin guardado en la nube ni correos enviados.")

        except Exception as e:
            log_queue.put(("error", f"❌ Error: {e}", "ERROR"))
        finally:
            matar_proceso_excel()
            try: ppt_app.Quit(); matar_proceso("POWERPNT.EXE")
            except: pass
            try: word_app.Quit(); matar_proceso("winword.exe")
            except: pass
            if sys.platform == "win32": pythoncom.CoUninitialize()
            log_queue.put(("done", "", "INFO"))

    def _poll(self):
        try:
            while True:
                item = log_queue.get_nowait()
                t = item[0]
                
                if t in ("log", "error", "done"):
                    msg, lvl = item[1], item[2]
                    self.con_box.config(state="normal")
                    tag = lvl if lvl in ("ERROR", "WARNING", "SUCCESS") else "INFO"
                    self.con_box.insert("end", msg + "\n", tag)
                    self.con_box.see("end"); self.con_box.config(state="disabled")

                elif t.startswith("res_"):
                    msg, lvl = item[1], item[2]
                    self.res_box.config(state="normal")
                    self.res_box.insert("end", msg + "\n\n", "TEC" if "TEC" in msg else "ESP")
                    self.res_box.see("end"); self.res_box.config(state="disabled")
                    
                    if len(item) == 7:
                        nombre_last, mat_last, oferta_last, periodo_last = item[3], item[4], item[5], item[6]
                        self.lbl_last_name.config(text=nombre_last, fg=SUCCESS)
                        self.lbl_last_details.config(
                            text=f"📌 Matrícula: {mat_last}\n📖 Oferta: {oferta_last}"
                        )

                elif t == "metric":
                    msg = item[1]
                    if msg == "tec":
                        self.tot_tec += 1
                        self.lbl_tec.config(text=str(self.tot_tec))
                    elif msg == "esp":
                        self.tot_esp += 1
                        self.lbl_esp.config(text=str(self.tot_esp))
                    elif msg == "correo":
                        self.tot_correos += 1
                        self.lbl_correos.config(text=str(self.tot_correos))
                    elif msg == "omitido":
                        self.tot_omitidos += 1
                        self.lbl_omit.config(text=str(self.tot_omitidos))

                if len(item) > 2 and item[2] == "ERROR": 
                    self.errores += 1
                    self.lbl_err.config(text=str(self.errores))
                    
                if t == "done": 
                    self.lbl_last_name.config(text="Ejecución Finalizada", fg=ACCENT)
                    self._finish(True)
                if t == "error": 
                    self._finish(False)
        except queue.Empty: pass
        except Exception as e: print(f"Error en poll: {e}")
        self.after(100, self._poll)

    def _finish(self, exitoso):
        self.running = False
        self.btn.config(text="▶  INICIAR", bg=ACCENT, command=self._start)
        if exitoso:
            self.status_label.config(text="● Proceso finalizado", fg=SUCCESS)
        else:
            self.status_label.config(text="● Proceso detenido con errores", fg=ERROR_COLOR)

if __name__ == "__main__":
    for d in [RUTA_BASES_DATOS, RUTA_MODULOS, RUTA_SALIDA, RUTA_IMAGENES, RUTA_RECORTES, RUTA_PLANTILLAS_NUEVAS, RUTA_PLANTILLAS_VIEJAS]:
        d.mkdir(exist_ok=True, parents=True)
    App().mainloop()